import tkinter as tk
from tkinter import filedialog, colorchooser, messagebox, ttk, filedialog
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from enum import Enum
import platform
from tkinter.constants import *
import customtkinter as ctk
import tkinter
import os
import openai
from PIL import Image, ImageTk
import requests
import io
 
class VerticalScrolledFrame(ttk.Frame):
    def __init__(self, parent, *args, **kw):
        ttk.Frame.__init__(self, parent, *args, **kw)
 
        # Create a canvas object and a vertical scrollbar for scrolling it.
        vscrollbar = ttk.Scrollbar(self, orient=VERTICAL)
        vscrollbar.grid(row=0, column=1, sticky=N + S)
        self.canvas = tk.Canvas(self, bd=0, highlightthickness=0,
                                width=200, height=300,
                                yscrollcommand=vscrollbar.set)
        self.canvas.grid(row=0, column=0, sticky=N + S + E + W)
        vscrollbar.config(command=self.canvas.yview)
 
        # Reset the view
        self.canvas.xview_moveto(0)
        self.canvas.yview_moveto(0)
 
        # Create a frame inside the canvas which will be scrolled with it.
        self.interior = ttk.Frame(self.canvas)
        self.interior_id = self.canvas.create_window((0, 0), window=self.interior, anchor=NW)
 
        self.interior.bind('<Configure>', self._configure_interior)
        self.canvas.bind('<Configure>', self._configure_canvas)

    def _configure_interior(self, event):
        # Update the scrollbars to match the size of the inner frame.
        size = (self.interior.winfo_reqwidth(), self.interior.winfo_reqheight())
        self.canvas.config(scrollregion=(0, 0, size[0], size[1]))
        if self.interior.winfo_reqwidth() != self.canvas.winfo_width():
            # Update the canvas's width to fit the inner frame.
            self.canvas.config(width=self.interior.winfo_reqwidth())

    def _configure_canvas(self, event):
        if self.interior.winfo_reqwidth() != self.canvas.winfo_width():
            # Update the canvas's width to fit the inner frame.
            self.canvas.config(width=self.interior.winfo_reqwidth())

         
    def _configure_canvas(self, event):
        if self.interior.winfo_reqwidth() != self.canvas.winfo_width():
            # Update the inner frame's width to fill the canvas.
            self.canvas.itemconfigure(self.interior_id, width=self.canvas.winfo_width())
         

class ImagePosition(Enum):
    TOP_LEFT = 1
    TOP_RIGHT = 2
    BOTTOM_LEFT = 3
    BOTTOM_RIGHT = 4

class PresentationApp:
    
    # List of fields to include in the textboxes
    fields_to_include = ["company", "location", "telephone", "description"]
    chosen_colors = []

    def __init__(self, master, *args, **kwargs):
        self.frame = VerticalScrolledFrame(master)
        self.frame.pack(expand = True, fill = tk.BOTH)
        #self.app = tk.Tk()
        #self.app.title("PowerPoint Generator")
        #ctk.set_appearance_mode("dark")
        self.slides = []
        self.image_positions = []
        self.text_formatting = {
            "job_id": {"fill": RGBColor(255, 255, 255), "color": RGBColor(0, 0, 0), "size": Pt(46), "bold": True, "italic": False, "underline": True},
            "company": {"fill": RGBColor(255, 255, 255), "color": RGBColor(0, 0, 0), "size": Pt(32), "bold": True, "italic": False, "underline": False},
            "location": {"fill": RGBColor(255, 255, 255), "color": RGBColor(0, 0, 0), "size": Pt(18), "bold": False, "italic": False, "underline": False},
            "telephone": {"fill": RGBColor(255, 255, 255), "color": RGBColor(0, 0, 0), "size": Pt(18), "bold": False, "italic": False, "underline": False},
            "description": {"fill": RGBColor(255, 255, 255), "color": RGBColor(0, 0, 0), "size": Pt(24), "bold": False, "italic": False, "underline": False}
        }
        self.images = []
        self.image_urls = []
        self.background_color = RGBColor(255, 255, 255)
        self.apply_background_colors_list = []
        self.chosen_colors = []
        self.current_color_index = 0
        self.init_ui()

    def init_ui(self):
        self.main_frame = ttk.Frame(self.frame.interior)
        self.main_frame.grid(row=0, column=0, padx=10, pady=10, sticky="nsew")

        self.left_frame = ttk.Frame(self.main_frame)
        self.left_frame.grid(row=0, column=0, padx=10, pady=10, sticky="w")
        
        #self.style = ttk.Style(self.main_frame)
        #self.style.theme_use("clam")
        #self.style.configure("TButton", font=("algerian", 10, "bold"), foreground="blue", background="yellow")

        self.job_id_label = ttk.Label(self.left_frame, text="Enter Job ID:")
        self.job_id_label.grid(row=0, column=0, padx=5, pady=5, sticky="w")

        self.job_id_entry = ttk.Entry(self.left_frame, width=40)
        self.job_id_entry.insert(0, "")
        self.job_id_entry.grid(row=0, column=1, padx=5, pady=5)

        self.company_label = ttk.Label(self.left_frame, text="Enter Company Name:")
        self.company_label.grid(row=0, column=2, padx=5, pady=5, sticky="w")

        self.company_entry = ttk.Entry(self.left_frame, width=40)
        self.company_entry.insert(0, "")
        self.company_entry.grid(row=0, column=3, padx=5, pady=5)

        self.location_count_label = ttk.Label(self.left_frame, text="Enter Location:")
        self.location_count_label.grid(row=1, column=0, padx=5, pady=5, sticky="w")

        self.location_entry = ttk.Entry(self.left_frame, width=40)
        self.location_entry.insert(0, "")
        self.location_entry.grid(row=1, column=1, padx=5, pady=5)

        self.telephone_count_label = ttk.Label(self.left_frame, text="Enter Telephone:")
        self.telephone_count_label.grid(row=1, column=2, padx=5, pady=5, sticky="w")

        self.telephone_entry = ttk.Entry(self.left_frame, width=40)
        self.telephone_entry.insert(0, "")
        self.telephone_entry.grid(row=1, column=3, padx=5, pady=5)

        self.num_colors_label = ttk.Label(self.left_frame, text="Enter the number of colors:")
        self.num_colors_label.grid(row=4, column=0, padx=5, pady=5, sticky="w")

        self.num_colors_spinbox = ttk.Spinbox(self.left_frame, from_=2, to=10, wrap=True, width=5)
        self.num_colors_spinbox.grid(row=4, column=1, padx=5, pady=5)

        choose_text_color_button = ttk.Button(self.left_frame, text="Choose Textboxes Background Color", command=self.create_color_buttons)
        choose_text_color_button.grid(row=4, column=2, padx=5, pady=5, columnspan=2)

        self.color_box = tk.Canvas(self.left_frame, width=50, height=50)
        self.color_box.grid(row=6, column=0, padx=5, pady=5, columnspan=2)

        self.textbox_width_label = ttk.Label(self.left_frame, text="Textbox Width")
        self.textbox_width_label.grid(row=7, column=0, padx=5, pady=5, sticky="w")

        self.textbox_width_entry = ttk.Entry(self.left_frame)
        self.textbox_width_entry.insert(0, "4.0")  # Default width
        self.textbox_width_entry.grid(row=7, column=1, padx=5, pady=5)

        self.textbox_height_label = ttk.Label(self.left_frame, text="Textbox Height")
        self.textbox_height_label.grid(row=7, column=2, padx=5, pady=5, sticky="w")

        self.textbox_height_entry = ttk.Entry(self.left_frame)
        self.textbox_height_entry.insert(0, "1.0")  # Default height
        self.textbox_height_entry.grid(row=7, column=3, padx=5, pady=5)

        self.every_slide_label = ttk.Label(self.left_frame, text="For every advert select image and enter description")
        self.every_slide_label.grid(row=9, column=0, padx=5, pady=5, columnspan=2)

        self.description_label = ttk.Label(self.left_frame, text="Enter Description:")
        self.description_label.grid(row=10, column=0, padx=5, pady=5, sticky="w")

        self.description_entry = ttk.Entry(self.left_frame, width=40)
        self.description_entry.insert(0, "")
        self.description_entry.grid(row=10, column=1, padx=5, pady=5)

        add_image_button = ttk.Button(self.left_frame, text="Add Image to Slide", command=self.choose_and_add_image)
        add_image_button.grid(row=11, column=0, padx=5, pady=5, columnspan=2)

        self.image_position_label = ttk.Label(self.left_frame, text="Select Image Position on the Slide")
        self.image_position_label.grid(row=11, column=2, padx=5, pady=5, sticky="w")
        
        # Create a label and text entry for the prompt
        self.for_generated_label = ttk.Label(self.left_frame, text="You can also generate and add an image using open ai")
        self.for_generated_label.grid(row=12, column=0, padx=5, pady=5, columnspan=4)
        
        self.prompt_label = ttk.Label(self.left_frame, text="prompt: ")
        self.prompt_label.grid(row=13, column=0, padx=5, pady=5, sticky="w")

        self.prompt_entry = ttk.Entry(self.left_frame, width=40)
        self.prompt_entry.insert(0, "")
        self.prompt_entry.grid(row=13, column=1, padx=5, pady=5, columnspan = 3)
        
        self.image_type = ttk.OptionMenu(self.left_frame, tk.StringVar(), "Generated", "Imported")
        self.image_type.grid(row=13, column=4, padx=5, pady=5)

        # Create a label and dropdown for selecting the style
        style_label = ttk.Label(self.left_frame, text="Style: ")
        style_label.grid(row=14, column= 0, padx=5, pady=5, sticky='w')

        self.style_dropdown = ttk.OptionMenu(self.left_frame, tk.StringVar(), "Realistic", "Cartoon", "3D Illustration", "Flat Art")
        self.style_dropdown.grid(row=14, column=1, padx=5, pady=5)

        # Create a label and scale for choosing the number of images
        self.number_label = ttk.Label(self.left_frame, text="# Images")
        self.number_label.grid(row=14, column = 2, padx=5, pady=5)

        self.number_slider = ttk.Scale(self.left_frame, from_=1, to=10, orient="horizontal")
        self.number_slider.grid(row =14, column = 3, padx = 5, pady = 5, columnspan = 2)

        # Create a button for generating images
        generate_images_button = ttk.Button(self.left_frame, text="Generate Images", command=self.generate_images)
        generate_images_button.grid(row=15, column = 0, padx=5, pady = 5, sticky = 'w')

        # Create a Canvas for displaying generated images
        self.generated_image_canvas = tk.Canvas(self.left_frame, width=512, height=512)
        self.generated_image_canvas.grid(row = 15, column = 2, columnspan = 4, rowspan = 5)

        # List to store the generated images
        self.generated_images = []
        self.current_image_index = 0

        self.image_position = ttk.Combobox(self.left_frame, values=[pos.name for pos in ImagePosition])
        self.image_position.grid(row=11, column=3, padx=5, pady=5, sticky = 'w')

        create_slide_button = ttk.Button(self.left_frame, text="Create Slide", command=self.create_slide)
        create_slide_button.grid(row=16, column=0, padx=5, pady=5, sticky = 'w')

        self.slide_count_label = ttk.Label(self.left_frame, text="Slides Created: 0")
        self.slide_count_label.grid(row=18, column=0, padx=10, pady=10, columnspan=2)

        create_presentation_button = ttk.Button(self.left_frame, text="Create Presentation", command=self.create_presentation)
        create_presentation_button.grid(row=16, column=1, padx=5, pady=5, sticky = 'w')

        reset_slide_count_button = ttk.Button(self.left_frame, text="Delete Created Slides", command=self.reset_slide_count)
        reset_slide_count_button.grid(row=17, column=0, padx=5, pady=5, sticky = 'w')

        reset_data_button = ttk.Button(self.left_frame, text="Reset Data", command=self.reset_data)
        reset_data_button.grid(row=17, column=1, padx=3, pady=5, sticky = 'w')
        
    def generate_images(self):
        openai.api_key = os.getenv("OPENAI_API_KEY")
        user_prompt = self.prompt_entry.get("0.0", ttk.END)
        user_prompt += " in style: " + self.style_dropdown.cget("text")

        response = openai.Image.create(
            prompt=user_prompt,
            n=int(self.number_slider.get()),
            size="512x512"
        )

        self.image_urls = [data['url'] for data in response['data']]

        # Clear any previous images
        self.generated_images.clear()
        self.current_image_index = 0

        for url in self.image_urls:
            response = requests.get(url)
            image = Image.open(io.BytesIO(response.content))
            photo_image = ImageTk.PhotoImage(image)
            self.generated_images.append(photo_image)

        self.display_generated_image()

    def display_generated_image(self):
        if self.generated_images:
            image = self.generated_images[self.current_image_index]
            self.generated_image_canvas.image = image
            self.generated_image_canvas.create_image(0, 0, anchor="nw", image=image)
            self.current_image_index = (self.current_image_index + 1) % len(self.generated_images)
            self.generated_image_canvas.after(3000, self.display_generated_image)


    def create_color_buttons(self):
        num_colors = int(self.num_colors_spinbox.get())

        # If color picker buttons frame already exists, destroy it
        if hasattr(self, "color_buttons_frame"):
            self.color_buttons_frame.destroy()

        # Create a new frame for the color picker buttons
        self.color_buttons_frame = ttk.Frame(self.left_frame)
        self.color_buttons_frame.grid(row=6, column=0, padx=5, pady=5, columnspan=4)

        # Create new color picker buttons inside the new frame
        self.color_picker_buttons = []
        for i in range(num_colors):
            button = ttk.Button(self.color_buttons_frame, text=f"Choose Color {i + 1}", command=lambda idx=i: self.choose_fill_color(idx))
            button.pack(side=tk.LEFT, padx=5)
            self.color_picker_buttons.append(button)


    def choose_fill_color(self, idx):
        print(f"Choosing color for index {idx}")
        color = colorchooser.askcolor(title=f"Choose Color {idx + 1}")
        if color[1]:
            # Append the chosen color to the list
            self.chosen_colors.append(RGBColor(*[int(channel) for channel in color[0]]))
            # Set canvas background to preview color
            self.color_box.config(bg=color[1])
        #print(f"Chosen colors: {self.chosen_colors}")

    def update_slide_count_label(self):
        count = len(self.slides)
        self.slide_count_label.config(text=f"Slides: {count}")

    def set_slide_background_color(self, slide, color):
        if color:
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = color

    def add_textbox_to_slide(self, slide, x, y, width, height, textbox_fill_color):

        # Helper function to add a textbox to a slide.
        textbox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(width), Inches(height))
        text_frame = textbox.text_frame
        # Set textbox fill color 
        textbox.fill.solid()
        textbox.fill.fore_color.rgb = textbox_fill_color

        return textbox, text_frame

    def add_text_to_textboxes(self, slide, slide_data, text_frames_list):

        # Iterate through fields_to_include and text_frames_list
        for field, text_frame in zip(self.fields_to_include, text_frames_list):
            # Check if there is content to add
            if field in slide_data:
                content = slide_data[field]
            else:
            # Set default content if no content is available
                content = ""
            p = text_frame.add_paragraph()
            p.text = content
            formatting = self.text_formatting.get(field, {})
            p.font.color.rgb = formatting.get("color", RGBColor(0, 0, 0))
            p.font.size = formatting.get("size", Pt(18))     

        return

    def create_slide(self):
        slide_data = {
            "company": self.company_entry.get(),
            "location": self.location_entry.get(),
            "telephone": self.telephone_entry.get(),
            "job_id": self.job_id_entry.get(),
            "description": self.description_entry.get(),
            "textbox_width": float(self.textbox_width_entry.get()),
            "textbox_height": float(self.textbox_height_entry.get()),
        }
        
        if not all(slide_data.values()):
            messagebox.showerror("Error", "Please enter text fields (Job ID, Company, Location, Telephone, and Description) before creating a slide.")
            return
        
        num_text_boxes = int(self.num_colors_spinbox.get())
        
        if not self.chosen_colors:
            messagebox.showerror("Error", "Please choose the number of background colors before creating a slide.")
            return
       
        self.slides.append((slide_data, self.chosen_colors))
        self.update_slide_count_label()

        position = ImagePosition[self.image_position.get()]
        self.image_positions.append(position)

    def choose_and_add_image(self):
        if self.image_type.cget("text") == "Imported":
            file_path = filedialog.askopenfilename(filetypes=[("Image Files", "*.png *.jpg *.jpeg *.gif *.bmp *.tif")])
            if file_path:
                self.images.append(file_path)
        else: 
            image_url = self.image_urls[self.current_image_index]
            
            first_word_match = re.search(r"\w+", self.user_prompt)
            if first_word_match:
                first_word = first_word_match.group()
            else:
                first_word = "image"
            # Generate a unique filename based on the first word and index
            filename = f"{first_word}_{self.current_image_index}.png"
            file_path = os.path.join("images", filename)

            response = requests.get(image_url)
            with open(file_path, "wb") as file:
                file.write(response.content)

    def add_image_to_slide(self, slide, image_path, position):
        width = Inches(2)  # You can adjust the width as needed
        height = Inches(2)  # You can adjust the height as needed
        
        print("Image Path:", image_path)
        image_extension = image_path.split('.')[-1]
        print("Image Extension:", image_extension)

        if position == ImagePosition.TOP_LEFT:
            left = Inches(0)
            top = Inches(0)
        elif position == ImagePosition.TOP_RIGHT:
            left = Inches(10) - width
            top = Inches(0)
        elif position == ImagePosition.BOTTOM_LEFT:
            left = Inches(0)
            top = Inches(7.5) - height
        elif position == ImagePosition.BOTTOM_RIGHT:
            left = Inches(10) - width
            top = Inches(7.5) - height
        else:
            # If the position is not recognized, place the image at the top-right
            left = Inches(10) - width
            top = Inches(0)

        image = slide.shapes.add_picture(image_path, left, top, width, height)
        
    def add_centered_textbox(self, prs, slide, text):
        slide_width = prs.slide_width
        slide_height = prs.slide_height
        textbox_width = Inches(6)
        textbox_height = Inches(1)
        left = (slide_width - textbox_width) / 2
        top = (slide_height - textbox_height) / 2  # Centers the textbox vertically

        textbox = slide.shapes.add_textbox(left, top, textbox_width, textbox_height)
        text_frame = textbox.text_frame
        p = text_frame.paragraphs[0]
        p.text = text
        p.alignment = PP_ALIGN.CENTER
        # Set the font size to 48
        for run in p.runs:
            run.font.size = Pt(48)

    def create_presentation(self):
        if len(self.slides) == 0:
            messagebox.showerror("Error", "Please create the required number of slides")
            return

        if not hasattr(self, "choose_fill_color"):
            messagebox.showerror("Error", "Please select a text color first.")
            return

        prs = Presentation()

        # Calculate colors and textboxes
        num_colors = len(self.chosen_colors)

        # Create Job ID Slide
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        job_id = self.job_id_entry.get()
        self.add_centered_textbox(prs, slide, f"{job_id}")
    
        slide_color_index = 0
        textbox_color_index = 0

        for i, (slide_data, colors) in enumerate(self.slides):
            slide = prs.slides.add_slide(prs.slide_layouts[6])

            # Set slide background color (cycling through self.chosen_colors)
            slide_color_index = i % num_colors
            background_color = self.chosen_colors[slide_color_index]
            self.set_slide_background_color(slide, background_color)

            # Create a list of available colors for the inner loop (excluding the outer loop color)
            available_colors = [c for j, c in enumerate(colors) if j != slide_color_index]

            y = 1.5  # Reset y for the next slide
            text_frames_list = []

            num_textboxes = max(4, len(available_colors))

            for j in range(num_textboxes):

                # Set textbox background color (cycling through self.chosen_colors)
                textbox_color_index = j % len(available_colors)
                textbox_color = available_colors[textbox_color_index]
                textbox, text_frame = self.add_textbox_to_slide(slide, 1, y, slide_data["textbox_width"], slide_data["textbox_height"], textbox_color)

                # Append the text_frame to the list
                text_frames_list.append(text_frame)

                y += 1
                # Increment the color index for the next textbox
                textbox_color_index = (textbox_color_index + 1) % len(available_colors)
                
            # Take only the first 4 text frames
            first_4_text_frames = text_frames_list[:4]
            if first_4_text_frames:
                 # Add text to the text_frame
                 self.add_text_to_textboxes(slide, slide_data, first_4_text_frames)
            
            if self.images and i < len(self.images):
                position = self.image_positions[i]
                self.add_image_to_slide(slide, self.images[i], position)

            # Increment the slide color index for the next slide
            slide_color_index = (slide_color_index + 1) % num_colors
            y = 1.5  # Reset y for the next slide


        file_name = self.job_id_entry.get()
        file_path = filedialog.asksaveasfilename(initialfile=file_name, defaultextension=".pptx", filetypes=[("PowerPoint files", "*.pptx")])
        if file_path:
            slide_data["job_id"] = file_path
            prs.save(file_path)
            
        print("PowerPoint presentation generated successfully!")

        self.reset_data()


    def reset_slide_count(self):
        self.slides = []
        self.image_positions = []
        self.update_slide_count_label()

    def reset_data(self):
        self.images = []
        self.image_urls=[]
        self.image_positions = []
        self.update_slide_count_label()
        
        # Reset the text formatting dictionary
        self.text_formatting = {
            "job_id": {"fill": RGBColor(255, 255, 255), "color": RGBColor(0, 0, 0), "size": Pt(18), "bold": False, "italic": False, "underline": False},
            "company": {"fill": RGBColor(255, 255, 255), "color": RGBColor(0, 0, 0), "size": Pt(18), "bold": False, "italic": False, "underline": False},
            "location": {"fill": RGBColor(255, 255, 255), "color": RGBColor(0, 0, 0), "size": Pt(32), "bold": False, "italic": False, "underline": False},
            "telephone": {"fill": RGBColor(255, 255, 255), "color": RGBColor(0, 0, 0), "size": Pt(32), "bold": False, "italic": False, "underline": False},
            "description": {"fill": RGBColor(255, 255, 255), "color": RGBColor(0, 0, 0), "size": Pt(18), "bold": False, "italic": False, "underline": False}
        }
        
        # Reset other fields to default values
        self.job_id_entry.delete(0, 'end')
        self.job_id_entry.insert(0, "")
        self.company_entry.delete(0, 'end')
        self.company_entry.insert(0, "")
        self.location_entry.delete(0, 'end')
        self.location_entry.insert(0, "")
        self.telephone_entry.delete(0, 'end')
        self.telephone_entry.insert(0, "")
        self.description_entry.delete(0, 'end')
        self.description_entry.insert(0, "")
        self.textbox_width_entry.delete(0, 'end')
        self.textbox_width_entry.insert(0, "4.0")
        self.textbox_height_entry.delete(0, 'end')
        self.textbox_height_entry.insert(0, "1.0")
        
        # Reset the spinbox to the default value
        self.num_colors_spinbox.delete(0, 'end')
        self.num_colors_spinbox.insert(0, "2")
        
        # Reset preview
        self.color_box.delete(all)
        self.color_box.config(bg="Silver")

        
        # If color picker buttons exist, destroy them
        if hasattr(self, "color_picker_buttons"):
            for button in self.color_picker_buttons:
                button.destroy()
            delattr(self, "color_picker_buttons")

        # Reset the slide count
        self.reset_slide_count()

    def run(self):
        self.frame.interior.mainloop()

if __name__ == "__main__":
    root = tk.Tk()
    root.title("")

    # Simply set the theme
    root.tk.call("source", "C:\Users\ryanw\OneDrive\Desktop\AutoAd - Copy\GUI-App-Python-main\Auto_ppt\openaippt")
    root.tk.call("set_theme", "light")
    app = PresentationApp(root)
    app.run()
