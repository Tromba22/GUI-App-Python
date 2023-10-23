import tkinter as tk
from tkinter import filedialog, colorchooser
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from enum import Enum
from tkinter import ttk

class ImagePosition(Enum):
    TOP_LEFT = 1
    TOP_RIGHT = 2
    BOTTOM_LEFT = 3
    BOTTOM_RIGHT = 4

class PresentationApp:
    def __init__(self):
        self.app = tk.Tk()
        self.app.title("PowerPoint Generator")
        self.slides = []
        self.image_positions = []
        self.text_formatting = {
            "title": {"fill": RGBColor(255, 255, 255), "color": RGBColor(0, 0, 0), "size": Pt(32), "bold": False, "italic": False, "underline": False},
            "company": {"fill": RGBColor(255, 255, 255), "color": RGBColor(0, 0, 0), "size": Pt(18), "bold": False, "italic": False, "underline": False},
            "job_id": {"fill": RGBColor(255, 255, 255), "color": RGBColor(0, 0, 0), "size": Pt(18), "bold": False, "italic": False, "underline": False},
            "description": {"fill": RGBColor(255, 255, 255), "color": RGBColor(0, 0, 0), "size": Pt(18), "bold": False, "italic": False, "underline": False}
        }
        self.images = []
        self.background_color = RGBColor(255, 255, 255)
        self.init_ui()

    def init_ui(self):
        create_slide_button = tk.Button(self.app, text="Create Slide", command=self.create_slide)
        create_slide_button.pack()

        choose_text_color_button = tk.Button(self.app, text="Choose Text Color", command=self.choose_text_color)
        choose_text_color_button.pack()

        choose_background_color_button = tk.Button(self.app, text="Choose Background Color", command=self.choose_background_color)
        choose_background_color_button.pack()

        add_image_button = tk.Button(self.app, text="Add Image to Slide", command=self.choose_and_add_image)
        add_image_button.pack()

        create_presentation_button = tk.Button(self.app, text="Create Presentation", command=self.create_presentation)
        create_presentation_button.pack()

        reset_slide_count_button = tk.Button(self.app, text="Reset Slide Count", command=self.reset_slide_count)
        reset_slide_count_button.pack()

        reset_data_button = tk.Button(self.app, text="Reset Data", command=self.reset_data)
        reset_data_button.pack()

        self.title_entry = tk.Entry(self.app, width=40)
        self.title_entry.insert(0, "Job Title")
        self.title_entry.pack()

        self.company_entry = tk.Entry(self.app, width=40)
        self.company_entry.insert(0, "Company")
        self.company_entry.pack()

        self.job_id_entry = tk.Entry(self.app, width=40)
        self.job_id_entry.insert(0, "Job ID")
        self.job_id_entry.pack()

        self.description_entry = tk.Entry(self.app, width=40)
        self.description_entry.insert(0, "Description")
        self.description_entry.pack()

        self.slide_count_label = tk.Label(self.app, text="Slides: 0")
        self.slide_count_label.pack()

        self.image_position = ttk.Combobox(self.app, values=[pos.name for pos in ImagePosition])
        self.image_position.set("TOP_LEFT")
        self.image_position.pack()

        self.textbox_width_entry = tk.Entry(self.app)
        self.textbox_width_entry.insert(0, "8.0")  # Default width
        self.textbox_width_entry.pack()

        self.textbox_height_entry = tk.Entry(self.app)
        self.textbox_height_entry.insert(0, "2.0")  # Default height
        self.textbox_height_entry.pack()

        self.title_fill_button = tk.Button(self.app, text="Title Fill", command=lambda: self.choose_fill_color("title"))
        self.title_fill_button.pack()

        self.company_fill_button = tk.Button(self.app, text="Company Fill", command=lambda: self.choose_fill_color("company"))
        self.company_fill_button.pack()

        self.job_id_fill_button = tk.Button(self.app, text="Job ID Fill", command=lambda: self.choose_fill_color("job_id"))
        self.job_id_fill_button.pack()

        self.description_fill_button = tk.Button(self.app, text="Description Fill", command=lambda: self.choose_fill_color("description"))
        self.description_fill_button.pack()

    def create_slide(self):
        slide = {
            "title": self.title_entry.get(),
            "company": self.company_entry.get(),
            "job_id": self.job_id_entry.get(),
            "description": self.description_entry.get(),
            "textbox_width": float(self.textbox_width_entry.get()),
            "textbox_height": float(self.textbox_height_entry.get()),
        }
        self.slides.append(slide)
        self.update_slide_count_label()

        position = ImagePosition[self.image_position.get()]

        self.image_positions.append(position)

    def update_slide_count_label(self):
        count = len(self.slides)
        self.slide_count_label.config(text=f"Slides: {count}")

    def choose_text_color(self):
        color = colorchooser.askcolor(title="Choose Text Color")
        if color[1]:
            self.text_color = RGBColor(*[int(channel) for channel in color[0]])

    def choose_fill_color(self, text_zone):
        color = colorchooser.askcolor(title=f"Choose Fill Color for {text_zone}")
        if color[1]:
            self.text_formatting[text_zone]["fill"] = RGBColor(*[int(channel) for channel in color[0]])

    def choose_background_color(self):
        color = colorchooser.askcolor(title="Choose Background Color")
        if color[1]:
            self.background_color = RGBColor(*[int(channel) for channel in color[0]])

    def choose_and_add_image(self):
        file_path = filedialog.askopenfilename(filetypes=[("Image Files", "*.png *.jpg *.jpeg *.gif *.bmp *.tif")])
        if file_path:
            self.images.append(file_path)

    def add_image_to_slide(self, slide, image_path, position):
        width = Inches(2)  # You can adjust the width as needed
        height = Inches(2)  # You can adjust the height as needed

        if position == ImagePosition.TOP_LEFT:
            left = Inches(0)
            top = Inches(0)
        elif position == ImagePosition.TOP_RIGHT:
            left = Inches(10) - width
            top = Inches(0)
        elif position == ImagePosition.BOTTOM_LEFT:
            left = Inches(0)
            top = Inches(7.5) - height
        elif position == ImagePosition.BOTTOM_RIGHT:
            left = Inches(10) - width
            top = Inches(7.5) - height
        else:
            # If the position is not recognized, place the image at the top-right
            left = Inches(10) - width
            top = Inches(0)

        image = slide.shapes.add_picture(image_path, left, top, width, height)

    def create_presentation(self):
        if len(self.slides) == 0:
            return

        prs = Presentation("output.pptx")

        total_slides_slide = prs.slides.add_slide(prs.slide_layouts[5])
        total_slides_title = total_slides_slide.shapes.title
        total_slides_title.text = "Total Number of Slides"
        total_slides_title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        total_slides_count = total_slides_slide.shapes.add_textbox(Inches(3), Inches(2), Inches(4), Inches(1))
        total_slides_text_frame = total_slides_count.text_frame
        total_slides_p = total_slides_text_frame.add_paragraph()
        total_slides_p.text = str(len(self.slides))
        total_slides_p.font.size = Pt(36)
        total_slides_p.alignment = PP_ALIGN.CENTER
        total_slides_slide.background.fill.solid()
        total_slides_slide.background.fill.fore_color.rgb = self.background_color

        for i, slide_data in enumerate(self.slides):
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            title = slide.shapes.title
            title.text = slide_data["title"]
            title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            if self.text_formatting["title"]["fill"]:
                title.fill.solid()
                title.fill.fore_color.rgb = self.text_formatting["title"]["fill"]
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = self.background_color

            company = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(slide_data["textbox_width"]), Inches(slide_data["textbox_height"]))
            company_text_frame = company.text_frame
            company_p = company_text_frame.add_paragraph()
            company_p.text = slide_data["company"]
            company_p.font.color.rgb = self.text_formatting["company"]["color"]
            company_p.font.size = self.text_formatting["company"]["size"]
            if self.text_formatting["company"]["fill"]:
                company.fill.solid()
                company.fill.fore_color.rgb = self.text_formatting["company"]["fill"]

            job_id = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(slide_data["textbox_width"]), Inches(slide_data["textbox_height"]))
            job_id_text_frame = job_id.text_frame
            job_id_p = job_id_text_frame.add_paragraph()
            job_id_p.text = slide_data["job_id"]
            job_id_p.font.color.rgb = self.text_formatting["job_id"]["color"]
            job_id_p.font.size = self.text_formatting["job_id"]["size"]
            if self.text_formatting["job_id"]["fill"]:
                job_id.fill.solid()
                job_id.fill.fore_color.rgb = self.text_formatting["job_id"]["fill"]

            # Use the user-specified width and height for the description textbox
            description = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(slide_data["textbox_width"]), Inches(slide_data["textbox_height"]))
            description_text_frame = description.text_frame
            description_p = description_text_frame.add_paragraph()
            description_p.text = slide_data["description"]
            description_p.font.color.rgb = self.text_formatting["description"]["color"]
            description_p.font.size = self.text_formatting["description"]["size"]
            if self.text_formatting["description"]["fill"]:
                description.fill.solid()
                description.fill.fore_color.rgb = self.text_formatting["description"]["fill"]

            if self.images and i < len(self.images):
                position = self.image_positions[i]
                self.add_image_to_slide(slide, self.images[i], position)

        file_path = filedialog.asksaveasfilename(defaultextension=".pptx", filetypes=[("PowerPoint files", "*.pptx")])
        if file_path:
            prs.save(file_path)
            print("PowerPoint presentation generated successfully!")

        self.reset_slide_count()

    def reset_slide_count(self):
        self.slides = []
        self.image_positions = []
        self.update_slide_count_label()

    def reset_data(self):
        self.slides = []
        self.images = []
        self.image_positions = []
        self.update_slide_count_label()

    def run(self):
        self.app.mainloop()

app = PresentationApp()
app.run()
